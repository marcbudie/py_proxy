"""Script om de py_proxy PowerPoint presentatie te genereren."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree


# Kleurenpalet
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # donker navy
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # catppuccin blauw
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # catppuccin groen
ACCENT3    = RGBColor(0xF3, 0x8B, 0xA8)   # catppuccin rood/roze
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xD0, 0xDA)
MUTED      = RGBColor(0x6C, 0x70, 0x86)
YELLOW     = RGBColor(0xF9, 0xE2, 0xAF)


def set_bg(slide, color: RGBColor):
    """Achtergrondkleur van een slide instellen."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def bullet_slide(prs, title_text, bullets, subtitle=None):
    """Slide met titel en bullet-lijst."""
    slide_layout = prs.slide_layouts[6]  # blanco
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, DARK_BG)

    W = prs.slide_width
    H = prs.slide_height

    # Accentlijn bovenaan
    bar = add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)

    # Titel
    add_textbox(slide, title_text,
                Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.8),
                font_size=32, bold=True, color=ACCENT)

    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.95), W - Inches(1), Inches(0.4),
                    font_size=16, color=LIGHT_GRAY)

    # Bullets
    y = Inches(1.4) if subtitle else Inches(1.1)
    for bullet in bullets:
        icon, text = bullet if isinstance(bullet, tuple) else ("•", bullet)
        # Icoon
        add_textbox(slide, icon,
                    Inches(0.4), y, Inches(0.4), Inches(0.45),
                    font_size=18, color=ACCENT2)
        # Tekst
        add_textbox(slide, text,
                    Inches(0.85), y, W - Inches(1.3), Inches(0.45),
                    font_size=18, color=WHITE)
        y += Inches(0.52)

    return slide


def two_column_slide(prs, title_text, left_items, right_items,
                     left_title="", right_title=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, DARK_BG)

    W = prs.slide_width
    H = prs.slide_height
    col_w = W // 2 - Inches(0.6)

    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)
    add_textbox(slide, title_text,
                Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.7),
                font_size=32, bold=True, color=ACCENT)

    for col_idx, (col_title, items, x_off) in enumerate([
        (left_title,  left_items,  Inches(0.4)),
        (right_title, right_items, W // 2 + Inches(0.2)),
    ]):
        y = Inches(1.0)
        if col_title:
            add_textbox(slide, col_title, x_off, y, col_w, Inches(0.45),
                        font_size=20, bold=True, color=ACCENT2)
            y += Inches(0.5)

        for item in items:
            icon, text = item if isinstance(item, tuple) else ("▸", item)
            add_textbox(slide, icon, x_off, y, Inches(0.35), Inches(0.42),
                        font_size=16, color=ACCENT3)
            add_textbox(slide, text, x_off + Inches(0.35), y,
                        col_w - Inches(0.35), Inches(0.42),
                        font_size=16, color=WHITE)
            y += Inches(0.48)

    # Verticale scheidingslijn
    line_shape = slide.shapes.add_shape(1,
        W // 2 - Inches(0.05), Inches(0.95),
        Inches(0.04), H - Inches(1.1))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = MUTED
    line_shape.line.fill.background()

    return slide


def code_slide(prs, title_text, code_text, caption=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, DARK_BG)

    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT)
    add_textbox(slide, title_text,
                Inches(0.5), Inches(0.2), W - Inches(1), Inches(0.7),
                font_size=32, bold=True, color=ACCENT)

    # Code blok achtergrond
    code_bg = add_rect(slide, Inches(0.4), Inches(1.05),
                       W - Inches(0.8), H - Inches(1.6),
                       RGBColor(0x11, 0x11, 0x1B))
    code_bg.line.color.rgb = MUTED
    code_bg.line.width = Pt(1)

    add_textbox(slide, code_text,
                Inches(0.6), Inches(1.15),
                W - Inches(1.2), H - Inches(1.8),
                font_size=13, color=ACCENT2, wrap=True)

    if caption:
        add_textbox(slide, caption,
                    Inches(0.5), H - Inches(0.45),
                    W - Inches(1), Inches(0.4),
                    font_size=13, color=MUTED)

    return slide


# ── Presentatie aanmaken ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Slide 1: Titelslide ─────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide1, DARK_BG)

# Grote accentrechthoek links
add_rect(slide1, 0, 0, Inches(0.18), H, ACCENT)

# Titel
add_textbox(slide1, "py_proxy",
            Inches(0.6), Inches(1.8), Inches(8), Inches(1.5),
            font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# Ondertitel
add_textbox(slide1,
            "Pure TCP SNI-proxy — transparant HTTPS-verkeer routeren\nop basis van hostnaam, zonder TLS te termineren",
            Inches(0.6), Inches(3.4), Inches(9), Inches(1.2),
            font_size=22, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# Versie / datum
add_textbox(slide1, "april 2026",
            Inches(0.6), H - Inches(1.0), Inches(4), Inches(0.5),
            font_size=15, color=MUTED)

# Decoratieve stippen rechts
for i, (x, y, sz, col) in enumerate([
    (Inches(11.0), Inches(1.5), Inches(1.8), ACCENT),
    (Inches(11.8), Inches(3.2), Inches(2.5), RGBColor(0x31, 0x32, 0x44)),
    (Inches(10.2), Inches(4.8), Inches(1.2), RGBColor(0x45, 0x47, 0x5A)),
]):
    circ = slide1.shapes.add_shape(9,  # OVAL
        x, y, sz, sz)
    circ.fill.solid()
    circ.fill.fore_color.rgb = col
    circ.line.fill.background()


# ── Slide 2: Wat is py_proxy? ───────────────────────────────────────────────
bullet_slide(prs, "Wat is py_proxy?", [
    ("🔀", "Pure TCP SNI-proxy — TLS wordt NIET getermineerd"),
    ("🔍", "Leest de SNI-hostnaam uit de TLS ClientHello"),
    ("📦", "Eén Python-bestand: proxy.py — geen frameworks"),
    ("⚙️",  "Runtime configuratie via config.json (live herladen)"),
    ("🐳", "Draait als systemd-service of Podman-container"),
    ("🛡️",  "Admin UI op HTTPS + TOTP/OTP-authenticatie"),
])

# ── Slide 3: Architectuur ───────────────────────────────────────────────────
code_slide(prs, "Netwerkarchitectuur",
"""Internet
  │
  ├─ :443  ──► proxy:8444  ──► SNI router ──► backend per hostname
  │                               │
  │                               ├─ voorbeeld.nl  ──► 192.168.1.10:443
  │                               ├─ andere.nl     ──► 192.168.1.11:443
  │                               └─ [uitgeschakeld] ──► 503 pagina
  │
  ├─ :2222 ──► proxy:2222  ──► tcp_route ──► 192.168.2.76:22  (SSH)
  └─ :300  ──► proxy:3333  ──► tcp_route ──► 192.168.2.76:300 (ThinLinc)

Admin UI (poort 9443) ──► alleen via SNI-route proxy.budie.eu
                          (niet direct open in firewall)""",
"Verkeer wordt transparant doorgestuurd — het backend-certificaat blijft intact")

# ── Slide 4: TLS-terminatie ─────────────────────────────────────────────────
bullet_slide(prs, "TLS-terminatie voor HTTP-backends",
    [
        ("✅", 'Activeer met  "tls_terminate": true  per route'),
        ("🔐", "Proxy termineert TLS aan client-kant (wildcard of per-route cert)"),
        ("📡", "Stuurt gedecrypteerde HTTP door naar de backend"),
        ("🏷️",  "Host-header wordt herschreven naar het backend-adres"),
        ("🔗", "X-Forwarded-For header automatisch toegevoegd"),
        ("⏱️",  "Idle-timeout na 30 seconden (keep-alive sessies)"),
    ],
    subtitle='Instelling in Admin UI: checkbox "HTTP backend (TLS termineren)"')

# ── Slide 5: TCP-routes & auto-disable ─────────────────────────────────────
two_column_slide(prs, "TCP-routes & Auto-uitschakelen",
    left_title="TCP-routes",
    left_items=[
        ("▸", "Plain TCP zonder TLS (SSH, RDP, …)"),
        ("▸", "Sleutel = luisterpoort (string)"),
        ("▸", "Geen SNI-inspectie"),
        ("▸", "Uitgeschakeld → verbinding direct verbroken"),
        ("▸", "Toggle via UI, Telegram of API"),
    ],
    right_title="Auto-uitschakelen",
    right_items=[
        ("▸", '"auto_disable_minutes": N  per route'),
        ("▸", "0 = nooit automatisch uitschakelen"),
        ("▸", "Timer start bij inschakelen"),
        ("▸", "Controle elke 15 seconden"),
        ("▸", "Telegram-melding bij auto-uitschakelen"),
        ("▸", "UI toont oranje afteltimer"),
    ],
)

# ── Slide 6: Admin UI ───────────────────────────────────────────────────────
bullet_slide(prs, "Admin UI  —  https://<host>:9443/",
    [
        ("🔒", "HTTPS verplicht — proxy weigert te starten zonder cert"),
        ("🔑", "Authenticatie via TOTP (Google Authenticator) of legacy OTP"),
        ("📋", "Overzicht alle TLS- en TCP-routes met aan/uit-toggle"),
        ("➕", "Route toevoegen via formulier (hostname, host, poort, label)"),
        ("🗑️",  "Route verwijderen met bevestigingsdialoog"),
        ("🔄", "Wijzigingen direct actief — geen herstart nodig"),
        ("⏱️",  "Auto-uitschakelen instellen per route"),
    ])

# ── Slide 7: Authenticatie ──────────────────────────────────────────────────
two_column_slide(prs, "Authenticatie",
    left_title="TOTP (aanbevolen)",
    left_items=[
        ("▸", "Google Authenticator / Authy / Bitwarden"),
        ("▸", "20 bytes (160 bits) geheim — cryptografisch veilig"),
        ("▸", "Constante-tijd vergelijking (hmac.compare_digest)"),
        ("▸", "Replay-bescherming via gebruikte time-steps"),
        ("▸", "Telegram-melding bij elke inlog"),
        ("▸", "Setup via /totp-setup in de UI"),
    ],
    right_title="Legacy OTP (fallback)",
    right_items=[
        ("▸", "8-cijferige code via e-mail én Telegram"),
        ("▸", "Geldig 5 minuten, eenmalig bruikbaar"),
        ("▸", "Rate limit: 60 seconden tussen aanvragen"),
        ("▸", "Blokkering na 10 foutieve pogingen"),
        ("▸", "Sessiecookie: HttpOnly; Secure; SameSite=Strict"),
        ("▸", "Sliding expiry: 30 minuten inactiviteit"),
    ],
)

# ── Slide 8: Telegram-bot ───────────────────────────────────────────────────
bullet_slide(prs, "Telegram-bot",
    [
        ("💬", "/status  — uptime, verbindingen, toggle-knoppen per route"),
        ("📜", "/logs    — laatste 30 regels uit journald"),
        ("🔄", "/reload  — config herladen (geen herstart)"),
        ("🔁", "/restart — service herstarten via systemctl"),
        ("📋", "/cert    — vervaldatums van alle certificaten (🟢🟡🔴)"),
        ("🔔", "Proactieve meldingen: start, backend onbereikbaar, dagelijks rapport"),
        ("🔐", "initData HMAC-validatie + auth_date tijdvenster (24 uur)"),
    ],
    subtitle="Draait als asyncio-task naast de proxy — geen apart proces")

# ── Slide 9: Installatie ────────────────────────────────────────────────────
code_slide(prs, "Installatie & Deployen",
"""# Interactief (kies systemd of container)
sudo bash install.sh

# Altijd systemd
sudo bash install.sh --systemd

# Altijd Podman-container (beheerd door systemd)
sudo bash install.sh --container

# Config herladen (SIGHUP — geen herstart)
systemctl reload  py-proxy

# Logs volgen
journalctl -u py-proxy -f

# Direct starten vanuit checkout (ontwikkeling)
python3 proxy.py""",
"install.sh overschrijft nooit config.json — veilig opnieuw uitvoeren bij updates")

# ── Slide 10: Beveiliging ───────────────────────────────────────────────────
two_column_slide(prs, "Beveiliging",
    left_title="Systeem & Netwerk",
    left_items=[
        ("🛡", "Dedicated gebruiker pyproxy (geen login shell)"),
        ("🔒", "Admin UI alleen via SNI-route (niet direct open)"),
        ("📝", "config.json nooit in git (secrets)"),
        ("🔏", "TLS private key permissiecheck bij elke start"),
        ("⏱", "TLS handshake timeout: 10 seconden"),
        ("🚫", "Verbindingen zonder geldig SNI direct verbroken"),
    ],
    right_title="Applicatie",
    right_items=[
        ("🛡", "html.escape() op alle foutpagina's (XSS)"),
        ("✅", "JSON-validatie op alle API-endpoints"),
        ("🔑", "TOTP: secrets.token_bytes(20) — niet random"),
        ("🔒", "OTP: secrets.randbelow() — niet random"),
        ("📵", "Gevoelige data gefilterd in /logs (Telegram)"),
        ("🔐", "SMTP via SSL/TLS — nooit plaintext"),
    ],
)

# ── Slide 11: Config-overzicht ──────────────────────────────────────────────
code_slide(prs, "config.json — overzicht",
"""{
  "listen_host": "0.0.0.0",
  "listen_ports": [8444],
  "tls_routes": {
    "voorbeeld.nl": {"host": "192.168.1.10", "port": 443,
                     "name": "label", "enabled": true},
    "mijnsite.nl":  {"host": "192.168.1.10", "port": 80,
                     "name": "mijnsite", "enabled": true,
                     "tls_terminate": true,
                     "auto_disable_minutes": 60}
  },
  "tcp_routes": {
    "2222": {"host": "192.168.1.10", "port": 22,
             "name": "ssh", "enabled": true}
  },
  "tls_cert": "/pad/naar/wildcard.crt",
  "tls_key":  "/pad/naar/wildcard.key",
  "telegram": {"bot_token": "...", "allowed_chat_ids": [123456789]}
}""",
"Live herladen via SIGHUP (systemctl reload) — geen herstart nodig")

# ── Slide 12: Samenvatting ──────────────────────────────────────────────────
bullet_slide(prs, "Samenvatting",
    [
        ("🔀", "SNI-proxy: transparant routeren zonder TLS te breken"),
        ("⚡", "Asyncio: efficiënt en lichtgewicht — één Python-bestand"),
        ("🔧", "Live configuratie via config.json + SIGHUP"),
        ("🛡️",  "Veilig: TOTP, rate-limiting, sessie-management"),
        ("📱", "Telegram-bot + Admin UI voor eenvoudig beheer"),
        ("🐳", "Flexibel: systemd of Podman-container"),
        ("🔄", "Auto-uitschakelen, naam-groepering, TLS-terminatie"),
    ])


out_path = "/home/user/py_proxy/py_proxy_presentatie.pptx"
prs.save(out_path)
print(f"Opgeslagen: {out_path}")
